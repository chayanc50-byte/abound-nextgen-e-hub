#!/usr/bin/env python3
"""Generate PowerPoint presentation from the video walkthrough script."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, bullets):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(6)

def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    # Add centered title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    add_title_slide(prs,
        "Abound Next-Gen E-Hub",
        "Website Walkthrough — Technical, Functional & Aesthetic Overview"
    )

    # 2. Intro
    add_content_slide(prs, "Introduction",
        [
            "E-commerce + multi-level business platform",
            "Covers: technical architecture, core features, design",
            "Recording tips: 8–12 min, 1080p, slow deliberate pace",
        ]
    )

    # 3. Public Pages & Aesthetics
    add_section_slide(prs, "1. Public Pages & Aesthetics", "0:45 – 2:30")

    add_content_slide(prs, "Navigation & Header",
        [
            "Responsive layout: desktop nav vs mobile hamburger",
            "Logo, nav links: Home, Products, Categories, About, Contact, FAQ, Join Us",
            "Login, Register, cart icon",
            "Sidebar menu on mobile",
        ]
    )

    add_content_slide(prs, "Homepage Sections",
        [
            "Product categories — 6 icon cards (Household, Grocery, Cosmetics, etc.)",
            "Popular products — cards with image, price (₹)",
            "Feature cards — Trusted Products, Business Opportunity, Growing Network, Support",
            "Business CTA — 'Build Your Own Business'",
            "Final CTA — 'Start Your Journey Today'",
        ]
    )

    add_content_slide(prs, "Homepage Technical Notes",
        [
            "CSS Grid: categories-grid, products-grid, features-grid",
            "Alternating section backgrounds",
            "CTAs lead to register, catalog, key flows",
        ]
    )

    # 4. Product Catalog
    add_section_slide(prs, "2. Product Catalog", "2:30 – 3:15")

    add_content_slide(prs, "Product Catalog",
        [
            "List all products with category filter",
            "Product grid: image, name, price",
            "Click → login required (if not logged in)",
            "Tech: Flask + Jinja2, SQLAlchemy (Product, Category)",
        ]
    )

    # 5. About & Director
    add_section_slide(prs, "3. About & Director Message", "3:15 – 4:00")

    add_content_slide(prs, "About & Director",
        [
            "About page: company info, 'Message from Director' card",
            "Director page: profile image, message, signature",
            "Consistent hero layout, card-based design",
        ]
    )

    # 6. Contact Support Center
    add_section_slide(prs, "4. Contact — Support Center", "4:00 – 4:45")

    add_content_slide(prs, "Contact Page Structure",
        [
            "Hero — 'Contact Us'",
            "Contact info — Phone, Email, Address, Facebook",
            "Support options — Call (tap-to-call), WhatsApp, Email",
            "Office address + Google Maps embed",
            "Contact form — Name, Email, Subject, Message",
            "Quick links — Home, About, FAQ, Login",
        ]
    )

    add_content_slide(prs, "Contact Technical",
        [
            "POST form → Flask backend",
            "Email via SMTP (env configurable)",
            "Responsive grid for support cards",
        ]
    )

    # 7. FAQ
    add_content_slide(prs, "5. FAQ", [
        "Accordion layout for common questions",
        "Expand/collapse interaction",
    ])

    # 8. Registration
    add_section_slide(prs, "6. Registration Flow", "5:00 – 6:15")

    add_content_slide(prs, "Registration & Setup",
        [
            "Register: email + referral code → Setup page",
            "Setup: name, username, password",
            "Welcome email (if configured)",
            "Tech: Unique referral IDs (ABN + 5 chars), parent-child hierarchy",
        ]
    )

    # 9. User Dashboard
    add_section_slide(prs, "7. Login & User Dashboard", "6:15 – 7:30")

    add_content_slide(prs, "Sales User Features",
        [
            "Dashboard — orders, team count, commissions",
            "Products — browse, place order (quantity, shipping)",
            "My Team — hierarchical team / downline view",
            "My Commissions — records from referred orders",
            "Profile, logout",
        ]
    )

    add_content_slide(prs, "User Technical",
        [
            "Role-based access: user, admin, super_admin",
            "Multi-level referral structure",
            "Commission on paid orders",
        ]
    )

    # 10. Admin Panel
    add_section_slide(prs, "8. Admin Panel", "7:30 – 9:30")

    add_content_slide(prs, "Admin Features",
        [
            "Dashboard — stats, quick actions",
            "Users — list, create, edit, delete",
            "Products & Categories — CRUD, image upload",
            "Orders — status updates (Pending → Shipped → Delivered)",
            "Commissions — list per user",
            "Wallets — earnings, balance, withdrawn",
            "Activity Logs — admin audit trail",
        ]
    )

    add_content_slide(prs, "Admin Technical",
        [
            "require_admin, require_super_admin decorators",
            "Models: User, Product, Order, Commission, Wallet, ActivityLog",
        ]
    )

    # 11. Technical Summary
    add_section_slide(prs, "9. Technical Summary", "9:30 – 10:15")

    add_content_slide(prs, "Tech Stack & Features",
        [
            "Backend: Flask, Python",
            "Database: SQLite via SQLAlchemy",
            "Auth: Werkzeug hashing, session-based",
            "Email: SMTP (welcome, orders, contact form)",
            "Frontend: Jinja2, vanilla CSS (flexbox, grid)",
            "Features: Multi-level hierarchy, referral IDs, commissions, RBAC",
        ]
    )

    # 12. Responsive & Closing
    add_section_slide(prs, "10. Responsive & Closing", "10:15 – 10:45")

    add_content_slide(prs, "Responsive Design",
        [
            "Header → hamburger on mobile",
            "Grids reflow for mobile",
            "Touch targets sized for mobile",
        ]
    )

    # 13. Key URLs
    add_content_slide(prs, "Key URLs",
        [
            "Home: /  |  Catalog: /catalog  |  About: /about",
            "Director: /about/director-message  |  Contact: /contact",
            "FAQ: /faq  |  Register: /register  |  Login: /login",
            "Dashboard: /dashboard  |  Products: /products",
            "My Team: /my-team  |  Commissions: /my-commissions",
            "Admin: /admin  |  Users: /admin/users  |  Products: /admin-products",
        ]
    )

    # 14. Thank you
    add_title_slide(prs,
        "Thank You",
        "Abound Next-Gen E-Hub — Full Website Walkthrough"
    )

    # Save
    output_path = "docs/Abound_NextGen_E-Hub_Walkthrough.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    main()
